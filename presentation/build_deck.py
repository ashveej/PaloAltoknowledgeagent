"""Generate the panel presentation deck (.pptx) — light theme matching the app,
structured to the interview agenda, with each PM focus area as a labeled slide.
Opens in PowerPoint/Keynote and imports into Google Slides.

Run:  ../.venv/bin/python build_deck.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import re

# --- App theme (matches the dashboard UI) ---
BG    = RGBColor(0xF6, 0xF7, 0xF9)   # app --bg
CARD  = RGBColor(0xFF, 0xFF, 0xFF)
TXT   = RGBColor(0x1A, 0x1D, 0x24)
MUT   = RGBColor(0x6B, 0x72, 0x80)
LINE  = RGBColor(0xE7, 0xE9, 0xEF)
ORANGE  = RGBColor(0xFA, 0x58, 0x2D)
ORANGE_D = RGBColor(0xE0, 0x41, 0x0F)
HI = RGBColor(0x1E, 0x8E, 0x3E); HIB = RGBColor(0xE7, 0xF4, 0xEC)
MED = RGBColor(0xC7, 0x77, 0x00); MEDB = RGBColor(0xFB, 0xF0, 0xDC)
LO = RGBColor(0xD6, 0x31, 0x2B); LOB = RGBColor(0xFB, 0xE9, 0xE8)
SUB = RGBColor(0x33, 0x38, 0x42)
SLATE = RGBColor(0x33, 0x3A, 0x4A)   # deterministic-code steps

prs = Presentation()
prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]; SW = prs.slide_width

TOTAL = 16  # set after building


def slide(bg=BG):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid(); s.background.fill.fore_color.rgb = bg
    return s


def box(s, l, t, w, h, fill=None, line=None, lw=1.0, rounded=True):
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
                             Inches(l), Inches(t), Inches(w), Inches(h))
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(lw)
    shp.shadow.inherit = False
    return shp


def setrun(p, text, size, color, bold=False, italic=False):
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.color.rgb = color
    r.font.bold = bold; r.font.italic = italic; r.font.name = "Arial"
    return r


def rich(p, text, size, base=TXT):
    for seg in re.split(r"(\*\*.*?\*\*)", text):
        if not seg:
            continue
        if seg.startswith("**") and seg.endswith("**"):
            setrun(p, seg[2:-2], size, TXT, bold=True)
        else:
            setrun(p, seg, size, base)


def tbox(s, l, t, w, h, anchor=None):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    if anchor:
        tf.vertical_anchor = anchor
    return tf


def eyebrow(s, text):
    setrun(tbox(s, 0.9, 0.55, 11.5, 0.4).paragraphs[0], text.upper(), 12.5, ORANGE_D, bold=True)


def heading(s, text):
    box(s, 0.9, 1.0, 0.16, 0.5, fill=ORANGE)
    setrun(tbox(s, 1.2, 0.9, 11.3, 0.9).paragraphs[0], text, 27, TXT, bold=True)


def footer(s, n):
    setrun(tbox(s, 0.9, 7.0, 7, 0.4).paragraphs[0], "paloalto  ·  Knowledge Agent", 10, MUT, bold=True)
    p = tbox(s, 11.0, 7.0, 1.6, 0.4).paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
    setrun(p, f"{n} / {TOTAL}", 10, MUT)


def bullets(s, items, top=2.2, left=0.95, width=11.4, size=16, gap=11):
    tf = tbox(s, left, top, width, 4.3); first = True
    for it in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
        p.space_after = Pt(gap)
        setrun(p, "●  ", size, ORANGE, bold=True); rich(p, it, size)


def two_col(s, left_items, right_items, top=2.2, size=15):
    bullets(s, left_items, top=top, left=0.95, width=5.6, size=size)
    bullets(s, right_items, top=top, left=6.95, width=5.5, size=size)


def table(s, headers, rows, top=2.2, col_w=None, size=13):
    g = s.shapes.add_table(len(rows) + 1, len(headers), Inches(0.95), Inches(top),
                           Inches(11.4), Inches(0.5 + 0.42 * len(rows))).table
    if col_w:
        for i, w in enumerate(col_w):
            g.columns[i].width = Inches(w)
    for j, h in enumerate(headers):
        c = g.cell(0, j); c.fill.solid(); c.fill.fore_color.rgb = TXT
        setrun(c.text_frame.paragraphs[0], h, size, CARD, bold=True)
    for i, row in enumerate(rows, 1):
        for j, val in enumerate(row):
            c = g.cell(i, j); c.fill.solid()
            c.fill.fore_color.rgb = CARD if i % 2 else BG
            rich(c.text_frame.paragraphs[0], val, size, base=(ORANGE_D if j == 0 else TXT))
    return g


def flow_diagram(s, top=2.35):
    """Color-coded pipeline: orange = AI step, slate = deterministic code."""
    stages = [("1", "Question", "start"), ("2", "Classify", "ai"),
              ("3", "Find\nsources", "code"), ("4", "Draft\n+ cite", "ai"),
              ("5", "Check\n& Score", "code"), ("6", "Polish", "ai"),
              ("7", "Answer", "start")]
    n = len(stages); bw = 1.32; bh = 1.15; gap = (11.4 - bw * n) / (n - 1)
    x = 0.95; centers = []
    for num, label, kind in stages:
        fill = ORANGE if kind == "ai" else (SLATE if kind == "code" else CARD)
        tc = CARD if kind in ("ai", "code") else TXT
        ln = LINE if kind == "start" else None
        shp = box(s, x, top, bw, bh, fill=fill, line=ln)
        tf = shp.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        setrun(p, num, 10, (CARD if kind in ("ai", "code") else MUT), bold=True)
        for part in label.split("\n"):
            q = tf.add_paragraph(); q.alignment = PP_ALIGN.CENTER
            setrun(q, part, 12.5, tc, bold=True)
        centers.append(x + bw / 2)
        if num != "7":
            ap = tbox(s, x + bw, top + 0.32, gap, 0.5).paragraphs[0]
            ap.alignment = PP_ALIGN.CENTER; setrun(ap, "▶", 13, ORANGE, bold=True)
        x += bw + gap
    # legend
    box(s, 8.4, 1.78, 0.22, 0.22, fill=ORANGE)
    setrun(tbox(s, 8.68, 1.7, 1.7, 0.35).paragraphs[0], "AI step", 11, MUT, bold=True)
    box(s, 10.1, 1.78, 0.22, 0.22, fill=SLATE)
    setrun(tbox(s, 10.38, 1.7, 2.2, 0.35).paragraphs[0], "Deterministic code", 11, MUT, bold=True)
    # escalation branch from box 5 (Check & Score)
    cx = centers[4]
    setrun(tbox(s, cx - 0.5, top + bh - 0.02, 1.0, 0.4).paragraphs[0], "▼", 13, LO, bold=True)
    box(s, 1.6, top + 1.55, 10.1, 0.78, fill=LOB, line=LO)
    rich(tbox(s, 1.85, top + 1.62, 9.6, 0.64, anchor=MSO_ANCHOR.MIDDLE).paragraphs[0],
         "Can't prove it?  →  route to the right **SME**  →  SME answers  →  "
         "**published as an approved source** (reusable next time)", 13)


def divider(s, n_of, title, time, focus):
    box(s, 0, 0, 13.333, 0.12, fill=ORANGE, rounded=False)
    setrun(tbox(s, 0.95, 2.4, 8, 0.5).paragraphs[0], f"SECTION {n_of}", 14, ORANGE_D, bold=True)
    box(s, 0.95, 3.05, 0.18, 0.95, fill=ORANGE)
    setrun(tbox(s, 1.3, 2.9, 11, 1.2).paragraphs[0], title, 40, TXT, bold=True)
    setrun(tbox(s, 1.3, 4.25, 11, 0.6).paragraphs[0], time, 18, MUT)
    if focus:
        b = box(s, 1.3, 5.1, 10.5, 0.9, fill=CARD, line=LINE)
        rich(tbox(s, 1.55, 5.28, 10.0, 0.6, anchor=MSO_ANCHOR.MIDDLE).paragraphs[0],
             "**Focus areas:** " + focus, 14)


# =================================================================== SLIDES
# 1 TITLE
s = slide()
box(s, 0, 0, 13.333, 0.12, fill=ORANGE, rounded=False)
b = box(s, 0.9, 1.5, 3.5, 0.5, fill=CARD, line=ORANGE)
setrun(b.text_frame.paragraphs[0], "TECHNICAL PM PANEL · 60 MIN", 12, ORANGE_D, bold=True)
tf = tbox(s, 0.9, 2.35, 11, 2.1)
setrun(tf.paragraphs[0], "Evidence-Governed Knowledge Agent", 42, TXT, bold=True)
p = tf.add_paragraph(); setrun(p, "for Technical Sales", 42, TXT, bold=True)
rich(tbox(s, 0.9, 4.55, 11.2, 1.5).paragraphs[0],
     "An evidence-governed RAG product — not a chatbot. Every answer is grounded in approved "
     "sources, cites the exact text, and carries a **deterministically-calculated** confidence "
     "score. Weak evidence routes to an SME, whose answer becomes reusable knowledge.", 18, SUB)
setrun(tbox(s, 0.9, 6.45, 11, 0.5).paragraphs[0],
       "Built spec-driven (Claude Code + Markdown specs) · FastAPI · LanceDB · Claude Opus 4.8 / Haiku 4.5", 13, MUT)

# 2 AGENDA
s = slide(); eyebrow(s, "Agenda"); heading(s, "How I'll use the 60 minutes")
table(s, ["Section", "Time", "What I'll cover"], [
    ["1 · Introduction", "5 min", "My background & Technical-PM philosophy"],
    ["2 · Context & Problem", "5 min", "The task (RAG) and the core user pain"],
    ["3 · Technical Deep Dive", "15 min", "Data strategy · system design & trade-offs · human-in-the-loop"],
    ["4 · Evaluation & Results", "10 min", "Metrics tied to KPIs · hard truths · retrospective"],
    ["5 · Q&A", "25 min", "Your questions"],
], top=2.3, col_w=[3.3, 1.3, 6.8], size=14); footer(s, 2)

# 3 DIVIDER — Introduction
s = slide(); divider(s, "1 of 5", "Introduction", "5 minutes",
                     "My background · my Technical-PM philosophy")

# 4 BACKGROUND + PHILOSOPHY
s = slide(); eyebrow(s, "Introduction"); heading(s, "My background & Technical-PM philosophy")
two_col(s,
    ["**Staff Product Manager — 15+ years** in enterprise CRM transformation & AI-enabled workflow automation (Walmart, DoorDash, Amplitude, Gusto, Akamai).",
     "**Applied AI track record:** a GenAI content platform with LLM orchestration; Salesforce Agentforce (intelligent routing, case deflection, AI summaries); predictive account scoring with Data Science.",
     "This product is my world rebuilt on an LLM — **knowledge, routing, approvals, and scoring**, grounded and governed."],
    ["**My operating system as a Technical PM:**",
     "**Slice the risk** — ship the thinnest vertical that retires the riskiest assumption first (here: *trust*).",
     "**Determinism where trust lives** — the LLM drafts language; code makes the trust decisions.",
     "**Make trust legible** — every score is explained; every claim shows its source."],
    top=2.15, size=14); footer(s, 4)

# 5 DIVIDER — Context & Problem
s = slide(); divider(s, "2 of 5", "Context & Problem", "5 minutes",
                     "The task (RAG) · the core user pain point")

# 6 PROBLEM + TASK
s = slide(); eyebrow(s, "Context & Problem"); heading(s, "The need, and the user pain it solves")
two_col(s,
    ["**The user:** Technical Sales team members who need answers about a **specific Palo Alto product**.",
     "**The need:** a system where they can ask a product question and get a **trustworthy, sourced** answer.",
     "**The pain:** product truth is **scattered** across docs, release notes, and wikis — and it **goes stale**.",
     "Experts re-answer the same questions; that knowledge is never captured or reused."],
    ["**The task = Retrieval-Augmented Generation (RAG)**, with a human-in-the-loop.",
     "RAG = retrieve approved documents first, then have the AI answer **only** from them, with citations.",
     "Not a chatbot: a generic chatbot hallucinates with no proof — a liability for a security vendor.",
     "_Analogy: give the team the official binder and say 'only quote from this, cite the page.'_"],
    top=2.2, size=15); footer(s, 6)

# 7 DIVIDER — Technical Deep Dive
s = slide(); divider(s, "3 of 5", "Technical Deep Dive", "15 minutes",
                     "Data Strategy · System Design & Trade-offs · The Human in the Loop")

# 8 ARCHITECTURE (visual flow)
s = slide(); eyebrow(s, "Technical Deep Dive · architecture"); heading(s, "End-to-end architecture")
flow_diagram(s, top=2.45)
b = box(s, 0.95, 5.95, 11.4, 0.95, fill=MEDB, line=MED)
rich(tbox(s, 1.25, 6.05, 10.9, 0.78, anchor=MSO_ANCHOR.MIDDLE).paragraphs[0],
     "**The key idea:** the **AI** (orange) only handles language — classify, draft, polish. Every **trust** "
     "decision — check citations & calculate confidence (slate) — is auditable code. That's why it's a "
     "governed system, not a chatbot.", 13.5)
footer(s, 8)

# 9 DATA STRATEGY
s = slide(); eyebrow(s, "Technical Deep Dive · Data Strategy"); heading(s, "Data: acquisition, synthesis & quality")
two_col(s,
    ["**Acquisition — three sources, on purpose:**",
     "**Real docs** — crawled docs.paloaltonetworks.com (Prisma Access, Cortex XDR, NGFW). ~17 pages, 68 passages.",
     "**PDFs** — ingested via pypdf + metadata sidecar.",
     "**Synthetic** — hand-generated stale + deprecated edge cases (can't test a freshness rule without a stale doc)."],
    ["**Quality & alignment:**",
     "Every passage carries product, version, source type, owner, **review date**, **deprecated flag**, approval status, votes.",
     "Metadata **drives** retrieval, scoring, and routing — deprecated never surfaces.",
     "**Honest gap:** the Cortex portal is JS-rendered → couldn't crawl it → I named the gap, not faked it."],
    top=2.1, size=14); footer(s, 9)

# 10 SYSTEM DESIGN & TRADE-OFFS
s = slide(); eyebrow(s, "Technical Deep Dive · System Design & Trade-offs"); heading(s, "Why these choices (build vs buy, latency vs accuracy)")
table(s, ["Decision", "Chose", "Over", "Why"], [
    ["Confidence", "Deterministic rubric", "LLM-judged", "Auditable; can't be gamed by fluency"],
    ["Grounding", "Hard code gates", "Prompt 'please cite'", "A gate guarantees; a prompt hopes"],
    ["Vector store", "LanceDB (embedded)", "Pinecone / pgvector", "Build vs buy: zero infra for the MVP"],
    ["Facts model", "Opus 4.8", "Haiku", "Accuracy where hallucination hurts"],
    ["Classifier", "Haiku 4.5", "Opus", "Latency / cost on a routing call"],
    ["Embeddings", "Local", "Voyage / OpenAI", "Free/offline MVP; one-line swap at scale"],
], top=2.0, col_w=[1.9, 2.5, 2.6, 4.4], size=12.5); footer(s, 10)

# 11 HUMAN IN THE LOOP
s = slide(); eyebrow(s, "Technical Deep Dive · The Human in the Loop"); heading(s, "The interface between AI logic & the user")
two_col(s,
    ["**For the Sales Engineer:**",
     "**Attribution UI** — the exact cited source text is shown inline; verify without leaving the screen.",
     "**Explainable confidence** — a pill + score breakdown + the cap reason.",
     "**Brand voice** — sounds consultative, but the Facts↔Style wall keeps facts unchanged."],
    ["**For the expert (SME):**",
     "**Escalate → route → email** the right SME group automatically.",
     "**Approve & Publish** → the answer becomes citable for everyone (the knowledge flywheel).",
     "**'This answer is wrong'** → lowers source reputation + opens a ticket (feedback → ranking, not retraining)."],
    top=2.1, size=14); footer(s, 11)

# 12 DIVIDER — Evaluation & Results
s = slide(); divider(s, "4 of 5", "Evaluation & Results", "10 minutes",
                     "Evaluation & Success (metrics tied to KPIs) · hard truths · Retrospective")

# 13 HOW I MEASURED SUCCESS
s = slide(); eyebrow(s, "Evaluation & Results"); heading(s, "How I measured success")
bullets(s, [
    "Not “does it sound smart?” — I set quality bars a salesperson can rely on:",
], top=1.9, size=15, gap=6)
table(s, ["What I measured", "Target"], [
    ["Every answer shows its source", "100%"],
    ["Answers drawn from outdated / retired documents", "0%"],
    ["A shaky answer shown as if it were solid", "0%  (it asks a human instead)"],
], top=2.5, col_w=[8.4, 3.0], size=14)
b = box(s, 0.95, 4.55, 11.4, 1.6, fill=HIB, line=HI)
tf = tbox(s, 1.25, 4.7, 10.9, 1.35, anchor=MSO_ANCHOR.MIDDLE)
setrun(tf.paragraphs[0], "How I validated it", 14, HI, bold=True)
p = tf.add_paragraph(); p.space_before = Pt(5)
rich(p, "I tested **15 real questions** — 5 easy, 5 medium, 5 hard. It **answered** the ones it could "
        "prove and **asked a human** for the rest. **15 / 15 correct behavior** — it never bluffed.", 15)
footer(s, 13)

# 14 HARD TRUTHS (plain language)
s = slide(); eyebrow(s, "Evaluation & Results"); heading(s, "What I learned — the hard truths")
bullets(s, [
    "**Don't trust the AI's own confidence.** Models sound certain even when wrong — so I have the system *calculate* confidence with rules, not ask the model.",
    "**Being too careful is also a failure.** Early on it sent answerable questions to a human. I tuned it to answer when it has proof — trust means answering when you can back it up, not just refusing.",
    "**Small rules need real-world testing.** A single “thumbs-down” could unfairly bury a good document. I only caught that by testing with real questions — and fixed it.",
], top=2.2, size=16, gap=16); footer(s, 14)

# 15 RETROSPECTIVE
s = slide(); eyebrow(s, "Evaluation & Results · Retrospective"); heading(s, "If I rebuilt it on today's stack")
two_col(s,
    ["**Would add / change:**",
     "**SME source certification** — experts verify & certify documents *before* they enter the agent, so every answer is built on accuracy-certified sources.",
     "**Agentic retrieval** — critique the evidence and re-search before answering; a tool-using agent for live lookups before escalating.",
     "Hosted embeddings + hybrid search + a feedback-tuned reranker; real Confluence/Jira connectors."],
    ["**Would NOT change — the spine:**",
     "Deterministic confidence.",
     "Hard citation gates.",
     "The Facts ↔ Style wall.",
     "_These earned the trust; everything else is an optimization._"],
    top=2.15, size=14); footer(s, 15)

# 16 Q&A
s = slide()
box(s, 0, 0, 13.333, 0.12, fill=ORANGE, rounded=False)
box(s, 0.95, 2.7, 0.18, 1.0, fill=ORANGE)
setrun(tbox(s, 1.3, 2.55, 11, 1.2).paragraphs[0], "Questions & Discussion", 42, TXT, bold=True)
rich(tbox(s, 1.3, 3.95, 10.8, 1.0).paragraphs[0],
     "Happy to go deeper on the architecture, the data strategy, the confidence rubric, the "
     "trade-offs, or the roadmap.", 18, SUB)
setrun(tbox(s, 1.3, 5.2, 11, 0.5).paragraphs[0],
       "Live demo available · repo + specs in the README", 13, MUT)
footer(s, 16)

prs.save("Executive_Presentation.pptx")
print(f"Saved Executive_Presentation.pptx — {len(prs.slides._sldIdLst)} slides (light app theme)")
