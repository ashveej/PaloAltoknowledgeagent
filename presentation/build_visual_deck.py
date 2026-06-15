"""Generate the VISUAL / non-technical deck (.pptx) — charts + plain language.
Opens in PowerPoint/Keynote and imports into Google Slides.

Run:  ../.venv/bin/python build_visual_deck.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
import re

WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SOFT  = RGBColor(0xF6, 0xF7, 0xF9)
TXT   = RGBColor(0x1A, 0x1D, 0x24)
MUT   = RGBColor(0x6B, 0x72, 0x80)
ORANGE  = RGBColor(0xFA, 0x58, 0x2D)
ORANGE_D = RGBColor(0xE0, 0x41, 0x0F)
LINE  = RGBColor(0xE7, 0xE9, 0xEF)
HI = RGBColor(0x1E, 0x8E, 0x3E); HIB = RGBColor(0xE7, 0xF4, 0xEC)
MED = RGBColor(0xC7, 0x77, 0x00); MEDB = RGBColor(0xFB, 0xF0, 0xDC)
LO = RGBColor(0xD6, 0x31, 0x2B); LOB = RGBColor(0xFB, 0xE9, 0xE8)
PRISMA = RGBColor(0xFA, 0x58, 0x2D); NGFW = RGBColor(0x2B, 0x6C, 0xB0); CORTEX = RGBColor(0x7C, 0x3A, 0xED)

prs = Presentation()
prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]; SW, SH = prs.slide_width, prs.slide_height


def slide(bg=WHITE):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid(); s.background.fill.fore_color.rgb = bg
    return s


def box(s, l, t, w, h, fill=None, line=None, line_w=1.0, rounded=True):
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
                             Inches(l), Inches(t), Inches(w), Inches(h))
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp


def setrun(p, text, size, color, bold=False, italic=False):
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.color.rgb = color; r.font.bold = bold
    r.font.italic = italic; r.font.name = "Arial"
    return r


def rich(p, text, size, base=TXT):
    for seg in re.split(r"(\*\*.*?\*\*)", text):
        if not seg:
            continue
        if seg.startswith("**") and seg.endswith("**"):
            setrun(p, seg[2:-2], size, TXT, bold=True)
        else:
            setrun(p, seg, size, base)


def tbox(s, l, t, w, h, anchor=None):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    if anchor:
        tf.vertical_anchor = anchor
    return tf


def eyebrow(s, text):
    setrun(tbox(s, 0.9, 0.55, 11.5, 0.4).paragraphs[0], text.upper(), 12.5, ORANGE_D, bold=True)


def heading(s, text):
    box(s, 0.9, 1.0, 0.16, 0.5, fill=ORANGE)
    setrun(tbox(s, 1.2, 0.9, 11.3, 0.9).paragraphs[0], text, 28, TXT, bold=True)


def card(s, l, t, w, h, icon, title, body):
    box(s, l, t, w, h, fill=SOFT, line=LINE)
    tf = tbox(s, l + 0.3, t + 0.25, w - 0.6, h - 0.5)
    setrun(tf.paragraphs[0], icon, 26, TXT)
    p = tf.add_paragraph(); p.space_before = Pt(4); setrun(p, title, 17, TXT, bold=True)
    p = tf.add_paragraph(); p.space_before = Pt(4); setrun(p, body, 13, MUT)


def pipeline(s, steps, top=2.3, branch=None):
    n = len(steps); gap = 0.18
    w = (11.5 - gap * (n - 1)) / n
    for i, (icon, title, sub) in enumerate(steps):
        l = 0.9 + i * (w + gap)
        box(s, l, top, w, 1.9, fill=WHITE, line=LINE)
        tf = tbox(s, l + 0.12, top + 0.18, w - 0.24, 1.6)
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER; setrun(p, icon, 24, TXT)
        p = tf.add_paragraph(); p.alignment = PP_ALIGN.CENTER; setrun(p, title, 13, TXT, bold=True)
        p = tf.add_paragraph(); p.alignment = PP_ALIGN.CENTER; setrun(p, sub, 10.5, MUT)
    if branch:
        box(s, 0.9, top + 2.2, 11.5, 0.95, fill=LOB, line=LO)
        tf = tbox(s, 1.2, top + 2.32, 11.0, 0.75, anchor=MSO_ANCHOR.MIDDLE)
        rich(tf.paragraphs[0], branch, 13.5)


# ----------------------------------------------------------------- 1 TITLE
s = slide(); box(s, 0, 0, 13.333, 0.12, fill=ORANGE, rounded=False)
b = box(s, 0.9, 1.5, 2.6, 0.5, fill=LOB, line=LO)
setrun(b.text_frame.paragraphs[0], "PRODUCT OVERVIEW", 12, ORANGE_D, bold=True)
tf = tbox(s, 0.9, 2.3, 11, 2.0)
setrun(tf.paragraphs[0], "Knowledge that Sales", 46, TXT, bold=True)
p = tf.add_paragraph(); setrun(p, "can actually trust.", 46, TXT, bold=True)
tf = tbox(s, 0.9, 4.6, 11.2, 1.5)
rich(tf.paragraphs[0],
     "An AI assistant for our sales teams that answers product questions **only from "
     "approved company sources** — and always shows its work and how sure it is.", 19, base=RGBColor(0x33, 0x38, 0x42))
setrun(tbox(s, 0.9, 6.4, 11, 0.5).paragraphs[0],
       "Everything in this deck is backed by the working system and real data.", 13, MUT)

# ----------------------------------------------------------------- 2 PROBLEM
s = slide(); eyebrow(s, "The problem"); heading(s, "Sales teams get hard questions — live, on the spot")
cw = 3.7
card(s, 0.9, 2.4, cw, 2.7, "⏱️", "Questions come fast",
     "“Does our product support X in version Y?” — mid-call with a customer.")
card(s, 4.82, 2.4, cw, 2.7, "\U0001F5C2️", "Scattered & goes stale",
     "Spread across manuals, release notes, and wikis — across products and versions. It ages.")
card(s, 8.74, 2.4, cw, 2.7, "\U0001F4C9", "A wrong answer costs the deal",
     "Confidently saying the wrong thing loses technical credibility — hard to win back.")
setrun(tbox(s, 0.9, 5.4, 11.5, 0.8).paragraphs[0],
       "A generic chatbot makes this worse — it sounds confident even when guessing, with no proof.", 18, RGBColor(0x33, 0x38, 0x42))

# ----------------------------------------------------------------- 3 IDEA
s = slide(); eyebrow(s, "The idea"); heading(s, "An expert assistant that never guesses")
tf = tbox(s, 0.9, 2.3, 7.0, 3.5)
setrun(tf.paragraphs[0], "Think of it as a diligent research assistant who:", 18, TXT, bold=True)
for t in ["\U0001F4D7  only quotes approved company sources",
          "\U0001F9FE  shows the exact quote behind every answer",
          "\U0001F3AF  tells you how confident it is — and why",
          "\U0001F64B  asks a human expert when it isn't sure"]:
    p = tf.add_paragraph(); p.space_before = Pt(10); setrun(p, t, 16, RGBColor(0x33, 0x38, 0x42))
b = box(s, 8.4, 2.4, 4.0, 3.0, fill=SOFT, line=LINE)
tf = tbox(s, 8.6, 3.0, 3.6, 2.2, anchor=MSO_ANCHOR.MIDDLE)
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER; setrun(p, "\U0001F6E1️", 40, TXT)
p = tf.add_paragraph(); p.alignment = PP_ALIGN.CENTER; setrun(p, "“No proof, no answer.”", 22, TXT, bold=True)
p = tf.add_paragraph(); p.alignment = PP_ALIGN.CENTER; setrun(p, "The one rule the whole system is built around.", 12, MUT)

# ----------------------------------------------------------------- 4 HOW IT WORKS
s = slide(); eyebrow(s, "How it works"); heading(s, "From question to trusted answer — in five steps")
pipeline(s, [
    ("\U0001F9ED", "1 · Understand", "Which product & topic?"),
    ("\U0001F4DA", "2 · Find sources", "Only approved, current docs."),
    ("✍️", "3 · Draft w/ quotes", "Only what sources say."),
    ("✅", "4 · Check & score", "Every claim cited; score confidence."),
    ("✨", "5 · Polish", "Sounds great — facts unchanged."),
], top=2.4,
   branch="\U0001F64B  **Not enough proof at step 4?**  The answer is blocked and sent to the right human expert instead. The system never bluffs.")

# ----------------------------------------------------------------- 5 SHOWS WORK
s = slide(); eyebrow(s, "It always shows its work · real example"); heading(s, "Every answer comes with the receipt")
box(s, 0.9, 2.3, 11.5, 0.7, fill=SOFT, line=LINE)
setrun(tbox(s, 1.1, 2.42, 11.1, 0.5).paragraphs[0],
       "❓  “What's the maximum number of BGP IP addresses per remote network?”", 15, TXT, bold=True)
b = box(s, 0.9, 3.15, 2.7, 0.5, fill=HIB, line=HI)
setrun(b.text_frame.paragraphs[0], "● High confidence · 97/100", 13, HI, bold=True)
tf = tbox(s, 0.9, 3.8, 11.4, 1.0)
rich(tf.paragraphs[0], "**Prisma Access supports a maximum of 250** local BGP IP addresses per "
     "remote network — 250 with a primary tunnel only, 125 with 2 links, or 62 with 4 links.", 16)
box(s, 0.9, 4.95, 11.4, 1.5, fill=SOFT, line=ORANGE, line_w=1.5)
tf = tbox(s, 1.15, 5.1, 11.0, 1.25)
setrun(tf.paragraphs[0], "\U0001F4D7 SOURCE: Prisma Access — “Onboard a Remote Network”", 11.5, MUT, bold=True)
p = tf.add_paragraph(); setrun(p, "“A maximum of 250 local BGP IP addresses is supported per remote "
       "network. This translates to 250 remote networks with a primary tunnel only…”", 13.5, RGBColor(0x26, 0x2B, 0x34))

# ----------------------------------------------------------------- 6 CONFIDENCE
s = slide(); eyebrow(s, "It tells you how sure it is"); heading(s, "Three confidence levels — calculated, not guessed")
levels = [("● High", HI, HIB, "Current, official sources fully back it. Use with confidence."),
          ("● Medium", MED, MEDB, "Backed, but the source is older or less authoritative. Double-check."),
          ("● Low", LO, LOB, "Not enough proof. Sent to a human expert instead of answered.")]
cw = 3.7
for i, (lab, col, bg, body) in enumerate(levels):
    l = 0.9 + i * (cw + 0.2)
    box(s, l, 2.3, cw, 1.9, fill=bg, line=col)
    tf = tbox(s, l + 0.25, 2.5, cw - 0.5, 1.6)
    setrun(tf.paragraphs[0], lab, 20, col, bold=True)
    p = tf.add_paragraph(); p.space_before = Pt(6); setrun(p, body, 13, TXT)
box(s, 0.9, 4.5, 11.5, 1.4, fill=SOFT, line=LINE)
tf = tbox(s, 1.15, 4.7, 11.0, 1.05, anchor=MSO_ANCHOR.MIDDLE)
rich(tf.paragraphs[0], "\U0001F4C5  **Real example:** a decryption answer was correct — but its source "
     "hadn't been reviewed in **241 days**, so the system automatically downgraded it from High to "
     "**Medium**. Freshness is built into the score.", 15)

# ----------------------------------------------------------------- 7 FLYWHEEL
s = slide(); eyebrow(s, "When it's not sure, a human steps in"); heading(s, "The knowledge flywheel")
pipeline(s, [
    ("\U0001F6AB", "Can't answer", "Not enough proof."),
    ("\U0001F4E7", "Email the expert", "Routed to the right SME."),
    ("\U0001F9D1‍\U0001F52C", "Expert replies", "Writes the authoritative answer."),
    ("\U0001F4BE", "Saved forever", "Becomes an approved source."),
    ("⚡", "Instant next time", "Everyone gets it immediately."),
], top=2.4)
setrun(tbox(s, 0.9, 5.0, 11.5, 0.9).paragraphs[0],
       "The system gets smarter every time an expert answers — yesterday's escalation is today's instant, cited answer.",
       18, RGBColor(0x33, 0x38, 0x42))

# ----------------------------------------------------------------- 8 DATA (donut)
s = slide(); eyebrow(s, "The knowledge behind it · real data"); heading(s, "What the assistant draws on today")
cd = CategoryChartData(); cd.categories = ["Prisma Access (8)", "NGFW (6)", "Cortex XDR (4)"]
cd.add_series("Sources", (8, 6, 4))
gf = s.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT, Inches(0.9), Inches(2.2), Inches(5.2), Inches(4.2), cd)
ch = gf.chart; ch.has_title = False
ch.has_legend = True; ch.legend.position = XL_LEGEND_POSITION.RIGHT; ch.legend.include_in_layout = False
ch.legend.font.size = Pt(13)
pts = ch.plots[0].series[0].points
for pt, col in zip(pts, [PRISMA, NGFW, CORTEX]):
    pt.format.fill.solid(); pt.format.fill.fore_color.rgb = col
for arr, (xl, yl, big, lab) in zip([0, 1], [(6.6, 2.4, "69", "searchable passages"), (9.6, 2.4, "3", "product lines")]):
    box(s, xl, yl, 2.7, 1.9, fill=SOFT, line=LINE)
    tf = tbox(s, xl, yl + 0.3, 2.7, 1.4); p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    setrun(p, big, 44, ORANGE_D, bold=True)
    p = tf.add_paragraph(); p.alignment = PP_ALIGN.CENTER; setrun(p, lab, 13, MUT)
box(s, 6.6, 4.5, 5.7, 1.9, fill=SOFT, line=LINE)
tf = tbox(s, 6.85, 4.75, 5.3, 1.45)
p = tf.paragraphs[0]
setrun(p, "14 fresh", 22, HI, bold=True); setrun(p, "  ·  ", 22, MUT); setrun(p, "4 flagged for review", 22, MED, bold=True)
p = tf.add_paragraph(); p.space_before = Pt(6)
setrun(p, "The system tracks how old every source is — and flags stale ones automatically.", 13, MUT)

# ----------------------------------------------------------------- 9 RESULTS
s = slide(); eyebrow(s, "Does it actually work? · real test"); heading(s, "We tested 15 real questions across difficulty")
rows = [("Easy questions", HI, "5 / 5 answered · High"),
        ("Medium questions", MED, "5 / 5 answered · Medium"),
        ("Hard / not in KB", LO, "5 / 5 correctly sent to an expert")]
for i, (lab, col, txt) in enumerate(rows):
    t = 2.4 + i * 0.85
    setrun(tbox(s, 0.9, t + 0.08, 2.6, 0.5).paragraphs[0], lab, 14, TXT, bold=True)
    b = box(s, 3.6, t, 8.7, 0.6, fill=col)
    tf = b.text_frame; tf.word_wrap = True; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
    setrun(p, txt + "   ", 13, WHITE, bold=True)
box(s, 0.9, 5.2, 11.4, 1.3, fill=HIB, line=HI)
tf = tbox(s, 1.2, 5.35, 11.0, 1.0, anchor=MSO_ANCHOR.MIDDLE)
p = tf.paragraphs[0]
setrun(p, "15 / 15   ", 30, HI, bold=True)
rich(p, "correct behavior — it **answered** what it could prove, and **escalated** the rest. It never bluffed.", 17)

# ----------------------------------------------------------------- 10 KPIs
s = slide(); eyebrow(s, "The numbers that matter"); heading(s, "How we define “good” — and hit it")
kpis = [("100%", "of answers carry a source citation"),
        ("0%", "use outdated / retired sources"),
        ("0", "low-confidence answers shown without an expert")]
cw = 3.7
for i, (big, lab) in enumerate(kpis):
    l = 0.9 + i * (cw + 0.2)
    box(s, l, 2.4, cw, 2.0, fill=SOFT, line=LINE)
    tf = tbox(s, l, 2.7, cw, 1.5); p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    setrun(p, big, 50, ORANGE_D, bold=True)
    p = tf.add_paragraph(); p.alignment = PP_ALIGN.CENTER; setrun(p, lab, 14, MUT)
setrun(tbox(s, 0.9, 4.8, 11.5, 1.2).paragraphs[0],
       "Success isn't “sounds smart.” It's every claim proven, nothing stale, and honesty about what it doesn't know.",
       18, RGBColor(0x33, 0x38, 0x42))

# ----------------------------------------------------------------- 11 VS CHATBOT
s = slide(); eyebrow(s, "Why this is different"); heading(s, "Not just another chatbot")
box(s, 0.9, 2.3, 5.6, 3.6, fill=WHITE, line=LO)
tf = tbox(s, 1.15, 2.5, 5.1, 3.2)
setrun(tf.paragraphs[0], "✕  A generic AI chatbot", 18, LO, bold=True)
for t in ["Answers from the open internet & memory", "Sounds confident even when wrong",
          "No proof, no source, no audit trail", "Never improves from our experts"]:
    p = tf.add_paragraph(); p.space_before = Pt(8); setrun(p, "•  " + t, 14, TXT)
box(s, 6.84, 2.3, 5.6, 3.6, fill=WHITE, line=HI)
tf = tbox(s, 7.1, 2.5, 5.1, 3.2)
setrun(tf.paragraphs[0], "✓  This knowledge assistant", 18, HI, bold=True)
for t in ["Answers only from approved sources", "Tells you how sure it is — and why",
          "Shows the exact quote every time", "Learns from every expert answer"]:
    p = tf.add_paragraph(); p.space_before = Pt(8); setrun(p, "•  " + t, 14, TXT)

# ----------------------------------------------------------------- 12 NEXT
s = slide(); eyebrow(s, "What's next"); heading(s, "Where this goes")
card(s, 0.9, 2.4, 3.7, 2.7, "\U0001F50C", "Connect more sources", "Wikis, ticketing, and live docs — automatically kept current.")
card(s, 4.82, 2.4, 3.7, 2.7, "\U0001F9D1‍\U0001F52C", "Help experts answer faster", "Pre-draft expert replies so they edit instead of write.")
card(s, 8.74, 2.4, 3.7, 2.7, "\U0001F4C8", "Measure deal impact", "Track faster answers, fewer interruptions, more wins.")
setrun(tbox(s, 0.9, 5.4, 11.5, 0.9).paragraphs[0],
       "The foundation is the part that earns trust — proof, honest confidence, and human experts in the loop.",
       18, RGBColor(0x33, 0x38, 0x42))

# ----------------------------------------------------------------- 13 CLOSE
s = slide(); box(s, 0, 0, 13.333, 0.12, fill=ORANGE, rounded=False)
b = box(s, 0.9, 2.0, 2.0, 0.5, fill=LOB, line=LO)
setrun(b.text_frame.paragraphs[0], "THANK YOU", 12, ORANGE_D, bold=True)
setrun(tbox(s, 0.9, 2.9, 11.3, 1.3).paragraphs[0], "Answers our teams can stand behind.", 36, TXT, bold=True)
rich(tbox(s, 0.9, 4.5, 11.2, 1.2).paragraphs[0],
     "Only approved sources · the exact proof shown · honest confidence · experts in the loop · smarter every day.",
     18, base=RGBColor(0x33, 0x38, 0x42))

prs.save("Executive_Presentation_Visual.pptx")
print(f"Saved Executive_Presentation_Visual.pptx — {len(prs.slides._sldIdLst)} slides")
